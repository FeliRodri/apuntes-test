import sys

try:
    import pptx
except ImportError:
    print("Falta la librería python-pptx. Por favor instala las dependencias:")
    print("pip install python-pptx")
    sys.exit(1)

def extract_text(filepath):
    """Extrae texto diapositiva por diapositiva de un PPTX."""
    text_content = []
    try:
        prs = pptx.Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                joined_text = "\n\n".join(slide_text)
                text_content.append(f"## Diapositiva {i+1}\n\n{joined_text}")
    except Exception as e:
        print(f"Error leyendo PPTX: {e}")
        sys.exit(1)
        
    return "\n\n".join(text_content)
